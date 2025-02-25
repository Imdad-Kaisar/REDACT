from abc import ABC, abstractmethod
from collections import defaultdict
import fitz
from docx import Document
import os
import csv
import json
import xml.etree.ElementTree as ET
from openpyxl import load_workbook
# Abstract Base Class for File Processors
class FileProcessor(ABC):
    def __init__(self, file_path):
        self.file_path = file_path

    @abstractmethod
    def extract_text(self):
        pass

    @abstractmethod
    def replace_text(self, replacement_map, output_path):
        pass

# PDF Processor
class PDFProcessor(FileProcessor):
    def extract_text(self):
        doc = fitz.open(self.file_path)
        text = ""
        for page in doc:
            text += page.get_text()
        return text
    def srgb_to_rgb(self , srgb_int):
      """Converts an sRGB color integer to RGB values.

      Args:
          srgb_int: The sRGB color as an integer.

      Returns:
          A tuple containing the RGB values (R, G, B) as integers
          in the range 0-255.
      """
    # Extract individual color components (R, G, B)
      blue = srgb_int & 0xFF
      green = (srgb_int >> 8) & 0xFF
      red = (srgb_int >> 16) & 0xFF
      blue = blue/255.0
      red = red/255.0 
      green = green/255.0
      return (red, green, blue)
    def replace_text(self, replacement_map, output_path="output.pdf"):
        doc = fitz.open(self.file_path)
        for page_index in range(len(doc)):
            page = doc[page_index]
            page.get_fonts()
            image_list = page.get_images()

            words = replacement_map.keys()
            redact = replacement_map
            hits = defaultdict()
            font_properties = defaultdict(dict)
            for word in words:
              hits[word] = page.search_for(word)  # list of rectangles where to replace
            blocks = page.get_text("dict", flags=11)["blocks"]
            for b in blocks:  # iterate through the text blocks
                for l in b["lines"]:  # iterate through the text lines
                    for s in l["spans"]:  # iterate through the text spans
                        for word in words:
                          if word in s["text"]:
                              for hit in hits[word]:
                                if hit[0] >= s["bbox"][0] and hit[1] >= s["bbox"][1] and hit[2] <= s["bbox"][2] and hit[3] <= s["bbox"][3]:
                                  font_properties[hit]["font"] = s.get("font", "helv")
                                  font_properties[hit]["size"] = s.get("size", 12)
                                  # span text color in sRGB format (int)
                                  font_properties[hit]["color"] = self.srgb_to_rgb(s.get("color", 0))
                                  print(font_properties[hit]["size"])
            for word in words:
              for rect in hits[word]:
                # Provide default values if font properties are missing
                font_name = font_properties.get(rect, {}).get("font", "tiro")
                font_size = font_properties.get(rect, {}).get("size", 12)
                font_color = font_properties.get(rect, {}).get("color", (0, 0, 0))
                font_name = "helv"
                p = 1
                rect[0] = rect[0]+p
                rect[1] = rect[1]+p
                rect[2] = rect[2]-p
                rect[3] = rect[3]-p
                # page.add_redact_annot(rect ,redact[word],fontname=font_name, fontsize=font_size,align=1,
                # text_color=font_color)
                page.add_redact_annot(rect, fill = (0,0,0))
                page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE) # Center horizontally
        doc.save(output_path, garbage=3, deflate=True)

# DOCX Processor
class DOCXProcessor(FileProcessor):
    def extract_text(self):
        doc = Document(self.file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text

    def replace_text(self, replacement_map, output_path="output.docx"):
        doc = Document(self.file_path)
        for paragraph in doc.paragraphs:
            for word, replacement in replacement_map.items():
                if word in paragraph.text:
                    paragraph.text = paragraph.text.replace(word, replacement)

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for word, replacement in replacement_map.items():
                        if word in cell.text:
                            cell.text = cell.text.replace(word, replacement)

        doc.save(output_path)

# Text File Processor (for TXT, CSV, XML)
class TextFileProcessor(FileProcessor):
    def extract_text(self):
        file_type = self._get_file_type()
        text = ""
        if file_type == "txt":
            with open(self.file_path, "r", encoding="utf-8") as file:
                text = file.read()
        # elif file_type == "csv":
        #     with open(self.file_path, "r", encoding="utf-8") as file:
        #         reader = csv.reader(file)
        #         text = "\n".join([", ".join(row) for row in reader])
        elif file_type == "xml":
            tree = ET.parse(self.file_path)
            root = tree.getroot()
            text = self._parse_xml_element(root)
        return text

    def replace_text(self, replacement_map, output_path="output"):
        file_type = self._get_file_type()
        if file_type == "txt":
            with open(self.file_path, "r", encoding="utf-8") as file:
                content = file.read()
            for word, replacement in replacement_map.items():
                content = content.replace(word, replacement)
            with open(output_path, "w", encoding="utf-8") as file:
                file.write(content)
        # elif file_type == "csv":
            # with open(self.file_path,"r",encoding="utf-8") as file:
            #   reader = csv.reader(file)
            #   rows = []
            #   for row in reader:
            #     new_row = [
            #       cell for cell in row
            #     ]
            #     for i , cell in enumerate(new_row):
            #         for word , replacement in replacement_map.items():
            #             if word in cell: 
            #               new_row[i] = cell.replace(word, replacement)
            #         rows.append(new_row)
            # with open(output_path , "w", encoding="utf-8",newline="") as file:
            #     writer = csv.writer(file)
            #     writer.writerows(rows)
        elif file_type == "xml":
            tree = ET.parse(self.file_path)
            root = tree.getroot()
            self._replace_xml_text(root, replacement_map)
            tree.write(output_path)

    def _get_file_type(self):
        _, ext = os.path.splitext(self.file_path)
        return ext.lower().lstrip(".")

    def _parse_xml_element(self, element):
        text = element.text.strip() if element.text else ""
        for child in element:
            text += self._parse_xml_element(child)
        return text

    def _replace_xml_text(self, element, replacement_map):
        if element.text:
            for word, replacement in replacement_map.items():
                element.text = element.text.replace(word, replacement)
        for child in element:
            self._replace_xml_text(child, replacement_map)

from pptx import Presentation
class PPTXProcessor(FileProcessor):
    def extract_text(self):
        """Extract all text from a PowerPoint file."""
        ppt = Presentation(self.file_path)
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text + "\n"
        return text

    def replace_text(self, replacement_map, output_path="output.pptx"):
        """Replace text in a PowerPoint file based on the replacement map."""
        ppt = Presentation(self.file_path)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for word, replacement in replacement_map.items():
                                if word in run.text:
                                    run.text = run.text.replace(word, replacement)
        ppt.save(output_path)
        print(f"Replaced text and saved to {output_path}")


class XLSXProcessor(FileProcessor):
    def extract_text(self):
        """Extract all text from an Excel file."""
        workbook = load_workbook(self.file_path)
        text = ""
        for sheet in workbook.sheetnames:
            worksheet = workbook[sheet]
            for row in worksheet.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        text += str(cell.value) + "\n"
        return text

    def replace_text(self, replacement_map, output_path="output.xlsx"):
        """Replace text in an Excel file based on the replacement map."""
        workbook = load_workbook(self.file_path)
        for sheet in workbook.sheetnames:
            worksheet = workbook[sheet]
            for row in worksheet.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        cell_value = str(cell.value)
                        for word, replacement in replacement_map.items():
                            if word in cell_value:
                                cell_value = cell_value.replace(word, replacement)
                        cell.value = cell_value
        workbook.save(output_path)
        print(f"Replaced text and saved to {output_path}")

import csv
import os
import requests

class CSVProcessor:
    def __init__(self, chunk_size=1000, temp_dir="temp_chunks", output_dir="processed_chunks"):
        self.chunk_size = chunk_size
        self.temp_dir = temp_dir
        self.output_dir = output_dir

    def split_csv(self, file_path):
        """Split a large CSV file into smaller chunks."""
        os.makedirs(self.temp_dir, exist_ok=True)
        file_count = 0

        with open(file_path, "r", encoding="utf-8") as file:
            print(file_path)
            reader = csv.reader(file)
            print(reader)
            headers = next(reader)
            
            if not headers:
                raise ValueError("The CSV file does not contain headers or is empty.")

            chunk = []
            for index, row in enumerate(reader, start=1):
                if not any(row):  # Skip empty rows
                    continue
                chunk.append(row)
                if index % self.chunk_size == 0:
                    output_file = os.path.join(self.temp_dir, f"chunk_{file_count}.csv")
                    self._write_csv(output_file, headers, chunk)
                    file_count += 1
                    chunk = []

            # Handle remaining rows
            if chunk:
                output_file = os.path.join(self.temp_dir, f"chunk_{file_count}.csv")
                self._write_csv(output_file, headers, chunk)
        return self.temp_dir

    def _write_csv(self, file_path, headers, rows):
        """Write a CSV file with the given headers and rows."""
        with open(file_path, "w", encoding="utf-8", newline="") as file:
            writer = csv.writer(file)
            # writer.writerow(headers)
            writer.writerows(rows)

    
    def process_chunk_remote(self, file_path, endpoint_url, gradation):
        """Send the chunk to an endpoint, get the processed data, and save it."""
        os.makedirs(self.output_dir, exist_ok=True)

        with open(file_path, "r", encoding="utf-8", newline="") as file:
            reader = csv.reader(file)
            text = "/n".join([", ".join(row) for row in reader])
            
        payload2 = {
            "text": text,
            "gradation_level": gradation
        }
        print(payload2)
        response = requests.post(endpoint_url, json=payload2, headers={"Content-Type": "application/json"})
        if response.status_code != 200:
            print(f"Failed to process chunk {file_path}: {response.text}")
            return
        
        replacement_map = response.json() if isinstance(response.json(), dict) else {} 
        print(replacement_map)
        with open(file_path,"r",encoding="utf-8") as file:
              reader = csv.reader(file)
              rows = []
              for row in reader:
                new_row = [
                  cell for cell in row
                ]
                for i , cell in enumerate(new_row):
                    for word , replacement in replacement_map.items():
                        if word in cell: 
                          new_row[i] = cell.replace(word, replacement)
                    rows.append(new_row)
        processed_file_path = os.path.join(self.output_dir, os.path.basename(file_path))

        with open(processed_file_path , "w", encoding="utf-8",newline="") as file:
            writer = csv.writer(file)
            writer.writerows(rows)
        
        # Save the processed chunk
        # processed_file_path = os.path.join(self.output_dir, os.path.basename(file_path))
        # with open(processed_file_path, "w", encoding="utf-8") as processed_file:
        #     processed_file.write(response)

    def merge_chunks(self, output_file):
        """Merge processed chunks into a single output CSV file."""
        chunk_files = sorted(
            [os.path.join(self.output_dir, file) for file in os.listdir(self.output_dir) if file.endswith(".csv")]
        )
        merged_rows = []

        for chunk_file in chunk_files:
            with open(chunk_file, "r", encoding="utf-8") as file:
                reader = csv.reader(file)
                headers = next(reader)
                merged_rows.extend(reader)
                
                # if headers is None:
                #     headers = next(reader)  # Read headers from the first chunk
                # else:
                #     next(reader, None)  # Skip headers for subsequent chunks
        print("hiiiiii---------------------")
        print(len(merged_rows))
        self._write_csv(output_file, headers, merged_rows)

    def cleanup(self):
        """Remove temporary and processed files."""
        import shutil
        shutil.rmtree(self.temp_dir, ignore_errors=True)
        shutil.rmtree(self.output_dir, ignore_errors=True)

# if _name_ == "_main_":
#     # Parameters
#     input_file = "large_file.csv"  # Path to your large CSV file
#     output_file = "processed_file.csv"  # Path to save the final merged output
#     endpoint_url = "https://example.com/api/process-csv"  # Your endpoint URL

#     # Initialize processor
#     processor = CSVProcessor(chunk_size=1000)

#     # Steps
#     processor.split_csv(input_file)

#     for chunk_file in os.listdir(processor.temp_dir):
#         processor.process_chunk_remote(os.path.join(processor.temp_dir, chunk_file), endpoint_url)

#     processor.merge_chunks(output_file)
#     processor.cleanup()

#     print(f"Processed file saved at: {output_file}")
# Update the DocumentProcessorFactory to support XLSX files
class DocumentProcessorFactory:
    @staticmethod
    def create_processor(file_path):
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            return PDFProcessor(file_path)
        elif ext == ".docx":
            return DOCXProcessor(file_path)
        elif ext == ".pptx":
            return PPTXProcessor(file_path)
        elif ext in [".txt",".csv", ".xml"]:
            return TextFileProcessor(file_path)
        elif ext == ".xlsx":
            return XLSXProcessor(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

# Factory Class
# class DocumentProcessorFactory:
#     @staticmethod
#     def create_processor(file_path):
#         ext = os.path.splitext(file_path)[1].lower()
#         if ext == ".pdf":
#             return PDFProcessor(file_path)
#         elif ext == ".docx":
#             return DOCXProcessor(file_path)
#         elif ext == ".pptx":
#             return PPTXProcessor(file_path)
#         elif ext in [".txt", ".csv", ".xml"]:
#             return TextFileProcessor(file_path)
#         else:
#             raise ValueError(f"Unsupported file type: {ext}")

# # Example Usage
# if __name__ == "__main__":
#     file_path = "test.pdf"  # Change to your file

#     replacement_map = {"Aryaman": "Yemen", "Name": "Dame"}

#     processor = DocumentProcessorFactory.create_processor(file_path)
#     print(processor.extract_text())

#     processor.replace_text(replacement_map)
